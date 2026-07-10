#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Génère la soutenance (Soutenance_Hischem_DiagPerf.pptx) — 17 diapos, figures intégrées."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

NAVY = RGBColor(0x1F, 0x2A, 0x44)
TEAL = RGBColor(0x1D, 0x9E, 0x75)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF1, 0xEF, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW = 13.333

def slide():
    return prs.slides.add_slide(BLANK)

def textbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def set_run(p, text, size, color=NAVY, bold=False, italic=False):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = "Calibri"
    return r

def title(s, text):
    tf = textbox(s, 0.6, 0.32, 12.1, 1.0)
    set_run(tf.paragraphs[0], text, 28, NAVY, bold=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.18), Inches(2.0), Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False

def bullets(s, items, l=0.8, t=1.6, w=11.7, h=5.4, size=20):
    tf = textbox(s, l, t, w, h)
    for i, it in enumerate(items):
        text, lvl = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(10); p.space_before = Pt(0)
        mark = "•  " if lvl == 0 else "–  "
        set_run(p, mark, size - lvl*2, TEAL if lvl == 0 else GREY)
        set_run(p, text, size - lvl*2, NAVY if lvl == 0 else GREY)

def fit(img, max_w, max_h):
    w, h = Image.open(img).size; ar = w / h
    iw, ih = max_w, max_w / ar
    if ih > max_h: ih, iw = max_h, max_h * ar
    return iw, ih

def footer(s, n):
    tf = textbox(s, 0.6, 7.02, 12.1, 0.35)
    p = tf.paragraphs[0]
    set_run(p, "De l'API au RAG — Hischem Hammoudi — DiagPerf", 9, GREY)
    pn = textbox(s, 12.3, 7.02, 0.8, 0.35).paragraphs[0]; pn.alignment = PP_ALIGN.RIGHT
    set_run(pn, str(n), 9, GREY)

def content(t, items, size=20):
    s = slide(); title(s, t); bullets(s, items, size=size); footer(s, len(prs.slides._sldIdLst))
    return s

def image_left(t, img, items):
    s = slide(); title(s, t)
    iw, ih = fit(img, 4.6, 5.0)
    s.shapes.add_picture(img, Inches(0.7 + (4.7 - iw) / 2), Inches(1.6 + (5.0 - ih) / 2), width=Inches(iw))
    bullets(s, items, l=5.7, t=1.7, w=6.9, size=18)
    footer(s, len(prs.slides._sldIdLst))
    return s

def image_center(t, img, caption=None, max_w=9.0, max_h=4.7):
    s = slide(); title(s, t)
    iw, ih = fit(img, max_w, max_h)
    s.shapes.add_picture(img, Inches((SW - iw) / 2), Inches(1.55), width=Inches(iw))
    if caption:
        tf = textbox(s, 1.0, 1.6 + ih + 0.1, 11.3, 0.8); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tf.paragraphs[0], caption, 15, GREY, italic=True)
    footer(s, len(prs.slides._sldIdLst))
    return s

def num(s, l, t, big, label, color=TEAL):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(3.4), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.fill.background(); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p, big, 30, color, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    set_run(p2, label, 14, NAVY)

def bug_slide(t, img, exchange, takeaway=None, max_w=11.8, max_h=4.5):
    s = slide(); title(s, t)
    iw, ih = fit(img, max_w, max_h)
    pic = s.shapes.add_picture(img, Inches((SW - iw) / 2), Inches(1.5), width=Inches(iw))
    pic.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); pic.line.width = Pt(0.75)
    tf = textbox(s, 0.8, 1.5 + ih + 0.15, 11.7, 0.8); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(tf.paragraphs[0], exchange, 16, NAVY, italic=True)
    if takeaway:
        tw = textbox(s, 0.8, 6.5, 11.7, 0.6); tw.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tw.paragraphs[0], takeaway, 16, TEAL, bold=True)
    footer(s, len(prs.slides._sldIdLst))
    return s

# ───────────────────────────── 1. Titre ─────────────────────────────
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(2.6))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background(); band.shadow.inherit = False
tf = textbox(s, 0.8, 0.7, 11.7, 1.6)
set_run(tf.paragraphs[0], "De l'API au RAG", 40, WHITE, bold=True)
p = tf.add_paragraph(); set_run(p, "Concevoir un assistant conversationnel fiable pour la relation client d'une PME automobile", 18, RGBColor(0xCF, 0xD6, 0xE4))
tf2 = textbox(s, 0.8, 3.0, 11.7, 3.0)
set_run(tf2.paragraphs[0], "Cas d'étude : Diagperf — automatisation du traitement des demandes clients", 18, NAVY, bold=True)
for txt in ["Hischem HAMMOUDI  ·  ING2 — Majeure Big Data & Machine Learning  ·  EFREI Paris",
            "Tuteur pédagogique : Antoine MILES  ·  Tuteur entreprise : Youcef ZAID (Diagperf)",
            "Soutenance — juin 2026"]:
    p = tf2.add_paragraph(); p.space_before = Pt(10); set_run(p, txt, 15, GREY)

# ───────────────────────────── 2. Contexte ─────────────────────────────
content("Contexte & enjeu", [
    "Diagperf : garage de reprogrammation moteur et de diagnostic électronique (Villenoy, équipe de 3)",
    "Demandes clients : devis, rendez-vous, questions techniques (E85, codes défauts) — surtout par téléphone, 1 à 4 / jour",
    "Traitement répétitif, mais à forte intensité de conseil (vérifier une compatibilité, annoncer un tarif exact)",
    "Enjeu : automatiser sans jamais produire d'information fausse",
    ("un prix ou une compatibilité erronés = perte de confiance, voire litige", 1),
])

# ───────────────────────────── 3. Problématique ─────────────────────────────
s = slide(); title(s, "Problématique & questions de recherche")
q = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.55), Inches(11.7), Inches(1.7))
q.fill.solid(); q.fill.fore_color.rgb = NAVY; q.line.fill.background(); q.shadow.inherit = False
qt = q.text_frame; qt.word_wrap = True; qt.vertical_anchor = MSO_ANCHOR.MIDDLE
qt.margin_left = Inches(0.3); qt.margin_right = Inches(0.3)
set_run(qt.paragraphs[0], "Comment concevoir un assistant conversationnel fiable pour une PME automobile, capable d'automatiser le traitement des demandes clients sans produire d'informations erronées ?", 18, WHITE, italic=True)
bullets(s, [
    "QR1 — Traiter des demandes variées en langage naturel, sans rupture de conversation ?",
    "QR2 — Le RAG réduit-il les hallucinations par rapport à un LLM par simple prompt ?",
    "QR3 — Où placer le savoir métier : figé dans le prompt, ou dans une base interrogeable ?",
], t=3.5, size=19)
footer(s, len(prs.slides._sldIdLst))

# ───────────────────────────── 4. Plan ─────────────────────────────
content("Plan", [
    "1.  État de l'art — les approches d'automatisation conversationnelle",
    "2.  Conception — l'assistant en trois temps : V0 → V1 → V2 (RAG)",
    "3.  Évaluation — un benchmark contrôlé V1 vs V2",
    "4.  Discussion, limites & perspectives",
], size=22)

# ───────────────────────────── 5. État de l'art ─────────────────────────────
content("État de l'art — 5 approches", [
    "Chatbots à règles / arbres de décision — déterministes, mais rigides (rupture de conversation)",
    "LLM par prompt + API — comprennent le langage, mais hallucinent et le savoir est figé",
    "RAG — ancrent chaque réponse dans une base documentaire interrogée à chaque requête",
    "Fine-tuning (LoRA) — internalise un style/comportement, pas les faits ; exige des données",
    "Agents / orchestration — action de bout en bout (perspective)",
], size=19)

# ───────────────────────────── 6. Pourquoi le RAG ─────────────────────────────
content("Pourquoi le RAG pour une PME", [
    "Dissocie la connaissance de la génération [Lewis et al., 2020]",
    "L'ancrage documentaire réduit fortement les hallucinations [Shuster et al., 2021]",
    "Maintenabilité : mettre à jour un tarif = éditer un fichier (aucun réentraînement)",
    "Adapté au contexte : peu de données → le fine-tuning est exclu",
], size=20)

# ───────────────────────────── 7. Pipeline ─────────────────────────────
image_left("La solution : un pipeline hybride", "figs/fig1.png", [
    "Webhook WhatsApp + signature HMAC",
    ("Flows déterministes pour les actions : devis, RDV, SAV", 1),
    ("RAG + Claude pour les questions libres", 1),
    "Sortie JSON typée : le LLM répond ou (re)bascule vers un flow",
    "→ souplesse du LLM + contrôle des actions sensibles",
])

# ───────────────────────────── 8. Vue client (captures) ─────────────────────────────
s = slide(); title(s, "Vue client : le parcours nominal")
welcome = "captures/Capture d'écran 2026-06-25 154905.png"
presta = "captures/Capture d'écran 2026-06-25 155003.png"
iw, ih = fit(welcome, 6.6, 3.5)
pw = s.shapes.add_picture(welcome, Inches(0.5), Inches(1.7), width=Inches(iw))
pw.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); pw.line.width = Pt(0.75)
bullets(s, [
    "Accueil : message de bienvenue + bouton « Nos prestations »",
    "Messages interactifs WhatsApp Business : boutons & listes natifs",
    "Parcours guidé — l'utilisateur n'a pas à tout taper",
], l=0.5, t=1.7 + ih + 0.3, w=7.6, size=16)
iw2, ih2 = fit(presta, 4.3, 4.7)
pp = s.shapes.add_picture(presta, Inches(8.9), Inches(1.55), width=Inches(iw2))
pp.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); pp.line.width = Pt(0.75)
tf = textbox(s, 8.5, 1.55 + ih2 + 0.1, 4.6, 0.6); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
set_run(tf.paragraphs[0], "Les 8 prestations proposées", 13, GREY, italic=True)
footer(s, len(prs.slides._sldIdLst))

# ───────────────────────────── 9. Évolution ─────────────────────────────
image_left("Une conception en trois temps", "figs/fig5.png", [
    "V0 — arbre de décision (boutons, sans LLM)",
    "V1 — LLM par prompt : comprend le langage… mais invente les faits",
    "V2 — RAG : réponses ancrées dans la base métier",
    "Chaque étape corrige une limite concrète de la précédente",
    ("(évolution datée par l'historique git)", 1),
])

# ─────────────── 9-11. Limites V1 — une capture par diapo ───────────────
bug_slide("V1 — rupture de conversation",
          "captures/Capture d'écran 2026-06-17 104936.png",
          "« Est-ce que je peux rouler tranquillement ? »  →  « Je n'ai pas bien saisi votre message »")
bug_slide("V1 — le flow à boutons rejette la question",
          "captures/Capture d'écran 2026-06-17 104959.png",
          "« Est-ce que je risque quelque chose après ma reprog ? »  →  « Merci de choisir un des stages »")
bug_slide("V1 — la question prise pour une plaque",
          "captures/Capture d'écran 2026-06-17 105331.png",
          "« Reprog pour un véhicule essence ? »  →  « Je n'ai pas reconnu la plaque »",
          takeaway="Trois ruptures distinctes → la motivation du passage au RAG.")

# ───────────────────────────── 10. Architecture RAG ─────────────────────────────
image_left("Architecture RAG (V2)", "figs/fig2.png", [
    "Indexation hors-ligne :",
    ("24 fichiers Q&A → embeddings Google (384 dim) → Supabase pgvector", 1),
    "Interrogation en ligne :",
    ("expansion synonymes → recherche hybride (vecteur + plein-texte) → rerank", 1),
    "Contexte ancré injecté dans le prompt de Claude Haiku 4.5",
])

# ─────────────── Côté technique : canal & backend ───────────────
content("Côté technique : WhatsApp Business + backend Supabase", [
    "Canal — WhatsApp Business Platform (Meta Cloud API)",
    ("messages interactifs (boutons, listes), texte et messages vocaux", 1),
    ("webhook entrant signé — HMAC-SHA256", 1),
    "Backend — Node.js + Express, hébergé sur Render (déploiement continu)",
    "Supabase (PostgreSQL) — double rôle :",
    ("état de conversation : où en est le client (machine à états)", 1),
    ("base de connaissances vectorielle (pgvector) pour le RAG", 1),
    "→ un seul backend orchestre : canal ↔ état ↔ RAG ↔ Claude",
], size=18)

# ───────────────────────────── 11. Protocole ─────────────────────────────
content("Le benchmark — protocole", [
    "31 questions : 28 ancrées sur la base + 3 cas réels de bugs",
    "4 conditions, croisant prompt (baké / dépouillé) × RAG (oui / non) :",
    ("V1 baké sans RAG · V2 baké + RAG · dépouillé sans RAG · dépouillé + RAG", 1),
    "Ablation : isole l'effet causal pur du RAG, toutes choses égales par ailleurs",
    "Juge : Claude Opus 4.8 — méthodologie « LLM-as-a-judge » [Zheng et al., 2023]",
], size=19)

# ───────────────────────────── 12. Résultats ─────────────────────────────
s = slide(); title(s, "Résultats")
iw, ih = fit("figs/fig_resultats.png", 7.6, 4.4)
s.shapes.add_picture("figs/fig_resultats.png", Inches(0.6), Inches(1.7), width=Inches(iw))
num(s, 9.0, 1.8, "0,66 → 0,98", "Exactitude (effet du RAG)")
num(s, 9.0, 3.5, "9 → 0", "Hallucinations", color=RED)
tf = textbox(s, 9.0, 5.3, 3.6, 1.4)
set_run(tf.paragraphs[0], "Le RAG quasi-double l'exactitude et supprime les hallucinations.", 15, NAVY, bold=True)
footer(s, len(prs.slides._sldIdLst))

# ───────────────────────────── 13. Le RAG ne règle pas tout ─────────────────────────────
content("Le RAG ne règle pas tout — 3 nuances", [
    "Baker le savoir dans le prompt peut tromper : règles qui se contaminent (9 hallucinations en V1)",
    "Le savoir baké entre en conflit avec le RAG (V2 : 1 hallucination résiduelle)",
    "Sur les codes défauts OBD : aucun apport (connaissance générale, déjà maîtrisée)",
    "En revanche, la rupture de conversation est bien corrigée",
], size=19)

# ───────────────────────────── 14. Recommandation ─────────────────────────────
s = slide(); title(s, "Recommandation d'ingénierie")
bullets(s, [
    "Sortir l'intégralité des faits (prix, compatibilités) du prompt → 100 % via le RAG",
    "La meilleure configuration n'est pas celle livrée, mais le prompt dépouillé + RAG :",
], t=1.6, size=20)
num(s, 1.3, 3.6, "0,98", "Exactitude")
num(s, 4.95, 3.6, "0", "Hallucination", color=RED)
num(s, 8.6, 3.6, "↓ coût", "Moins cher")
tf = textbox(s, 0.8, 5.5, 11.7, 1.0)
set_run(tf.paragraphs[0], "Plus exact, moins cher, et plus simple à maintenir (un tarif = un fichier édité).", 18, NAVY, bold=True)
footer(s, len(prs.slides._sldIdLst))

# ───────────────────────────── 15. Limites & perspectives ─────────────────────────────
content("Limites & perspectives", [
    "Limites :",
    ("benchmark hors-ligne (assistant pas encore déployé)", 1),
    ("juge LLM faillible ; qualité bornée par celle des sources", 1),
    "Perspectives :",
    ("agents outillés : devis et RDV automatiques de bout en bout", 1),
    ("extension multilingue (FR / AR) ; affinage du ton par fine-tuning léger", 1),
], size=19)

# ───────────────────────────── 16. Conclusion ─────────────────────────────
content("Conclusion & apports", [
    "Le RAG rend l'assistant fiable et maintenable pour une PME — démontré, chiffres à l'appui",
    "Apports techniques : pipeline RAG de bout en bout, déploiement, évaluation rigoureuse",
    "Apports transverses : conduite de projet en autonomie, de la V0 à la mise en production",
    "Compétences iziA couvertes : 2.1 analyser → 2.2 concevoir → 2.3 mettre en œuvre → 2.4 exploiter",
], size=19)

# ───────────────────────────── 17. Merci ─────────────────────────────
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(7.5))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background(); band.shadow.inherit = False
tf = textbox(s, 0.8, 2.8, 11.7, 2.0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p, "Merci de votre attention", 36, WHITE, bold=True)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(14)
set_run(p2, "Questions ?", 22, TEAL)

prs.save("Soutenance_Hischem_DiagPerf.pptx")
print("OK — diapos :", len(prs.slides._sldIdLst))
